"""Build MBR PPTX slides matching the Moxie design system."""

import os
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .data_schema import MBRData
from .charts import make_gauge_png

# ── Design tokens ──────────────────────────────────────────────
SLIDE_W = Inches(8.5)
SLIDE_H = Inches(11)

PURPLE = RGBColor(0x4D, 0x17, 0x51)
GOLD = RGBColor(0xB8, 0x93, 0x3A)
PINK_ACCENT = RGBColor(0xC2, 0x7B, 0xA0)
RED = RGBColor(0xB5, 0x36, 0x31)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x8C, 0x85, 0x80)
LIGHT_GRAY = RGBColor(0x91, 0x8A, 0x85)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_BG = RGBColor(0xFD, 0xF6, 0xEE)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER_BG = RGBColor(0x4D, 0x17, 0x51)
TRACK_COLOR = RGBColor(0xD2, 0xCB, 0xC4)

FONT_HEADING = "Playfair Display"
FONT_BODY = "Arial"
FONT_LABEL = "Arial"

ASSETS_DIR = Path(__file__).parent.parent / "assets"
COPYRIGHT = "Moxie Partners, Inc. Copyright \u00a92024. Private and confidential."


def _add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(12), font_color=PURPLE, bold=False, alignment=PP_ALIGN.LEFT,
                 word_wrap=True):
    """Add a textbox with a single run of formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def _add_card(slide, left, top, width, height, fill_color=CARD_BG, corner_radius=Inches(0.08)):
    """Add a rounded rectangle card shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust corner radius via XML
    from lxml import etree
    sp = shape._element
    prst = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prst is not None:
        avLst = prst.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prst, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(avLst):
                avLst.remove(child)
        gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', 'val 5000')
    return shape


def _add_header_bar(slide, practice_name, month_year, logo_path=None):
    """Add the header bar with practice name and optional logo."""
    # Practice name + date
    _add_textbox(slide, Inches(0.51), Inches(0.03), Inches(4.5), Inches(0.45),
                 f"{practice_name} | {month_year}",
                 font_name=FONT_HEADING, font_size=Pt(8), font_color=PURPLE, bold=True)

    # Moxie logo
    logo = logo_path or str(ASSETS_DIR / "moxie_logo.png")
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(7.0), Inches(0.05), height=Inches(0.35))


def _add_copyright(slide):
    """Add copyright footer."""
    _add_textbox(slide, Inches(0.5), Inches(10.55), Inches(7.5), Inches(0.3),
                 COPYRIGHT, font_size=Pt(6), font_color=GRAY, alignment=PP_ALIGN.CENTER)


def _fmt_dollar(val, show_cents=False):
    """Format a dollar value."""
    if val is None:
        return "N/A"
    if show_cents:
        return f"${val:,.2f}"
    if abs(val) >= 1000:
        return f"${val:,.0f}"
    return f"${val:,.2f}"


def _fmt_pct(val):
    """Format a percentage from decimal."""
    if val is None:
        return "N/A"
    p = val * 100
    if p == int(p):
        return f"{int(p)}%"
    return f"{p:.1f}%"


def _set_bg(slide, color=CREAM_BG):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide builders ─────────────────────────────────────────────

def build_slide_1_cover(prs, data: MBRData):
    """Slide 1: Cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_bg(slide)

    # Practice name (large, centered, purple)
    _add_textbox(slide, Inches(0.5), Inches(3.5), Inches(7.5), Inches(1.3),
                 data.practice_name,
                 font_name=FONT_HEADING, font_size=Pt(44), font_color=PURPLE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(6.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run1 = p.add_run()
    run1.text = "Monthly Business Review"
    run1.font.name = FONT_HEADING
    run1.font.size = Pt(24)
    run1.font.color.rgb = GOLD
    run1.font.bold = False

    from pptx.oxml.ns import qn
    import copy
    br = p.add_run()
    br.text = "\n"

    run2 = p.add_run()
    run2.text = f"{data.month_name} {data.year}"
    run2.font.name = FONT_HEADING
    run2.font.size = Pt(24)
    run2.font.color.rgb = GOLD

    # Logo top right
    logo = str(ASSETS_DIR / "moxie_logo.png")
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(6.5), Inches(0.4), height=Inches(0.5))

    _add_copyright(slide)


def build_slide_2_snapshot(prs, data: MBRData):
    """Slide 2: Monthly Snapshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    # Moxie start month
    if data.moxie_start_month:
        _add_textbox(slide, Inches(0.55), Inches(0.45), Inches(3.3), Inches(0.25),
                     f"Moxie Start Month: {data.moxie_start_month}",
                     font_size=Pt(7), font_color=GRAY)

    # "Monthly Snapshot" title
    _add_textbox(slide, Inches(0.56), Inches(0.75), Inches(3.5), Inches(0.4),
                 "Monthly Snapshot",
                 font_name=FONT_HEADING, font_size=Pt(16), font_color=PINK_ACCENT, bold=True)

    # ── Executive Summary ──
    _add_card(slide, Inches(0.38), Inches(1.10), Inches(7.89), Inches(1.16))
    if data.executive_summary:
        _add_textbox(slide, Inches(0.56), Inches(1.14), Inches(7.37), Inches(1.05),
                     data.executive_summary,
                     font_size=Pt(9), font_color=DARK_TEXT)

    # ── Key Metrics Cards (4 cards) ──
    card_w = Inches(1.85)
    card_h = Inches(1.52)
    card_y = Inches(2.45)
    card_starts = [Inches(0.38), Inches(2.35), Inches(4.32), Inches(6.29)]

    metrics = [
        ("Monthly Net Revenue", _fmt_dollar(data.monthly_net_revenue), data.revenue_mom_pct),
        ("Total Appointments", str(data.total_appointments), data.appointments_mom_pct),
        ("AOV", _fmt_dollar(data.aov), data.aov_mom_pct),
        ("Quarter to Date", _fmt_dollar(data.quarter_to_date) if data.quarter_to_date else "N/A", None),
    ]

    for i, (label, value, mom) in enumerate(metrics):
        x = card_starts[i]
        _add_card(slide, x, card_y, card_w, card_h)

        _add_textbox(slide, x + Inches(0.13), card_y + Inches(0.19), card_w - Inches(0.2), Inches(0.42),
                     label, font_name=FONT_HEADING, font_size=Pt(9.5), font_color=PURPLE, bold=True)

        _add_textbox(slide, x + Inches(0.18), card_y + Inches(0.54), card_w - Inches(0.2), Inches(0.43),
                     value, font_size=Pt(20), font_color=PURPLE, bold=True)

        if mom is not None:
            arrow = "\u2191" if mom >= 0 else "\u2193"
            color = GREEN if mom >= 0 else RED
            _add_textbox(slide, x + Inches(0.18), card_y + Inches(1.06), card_w - Inches(0.2), Inches(0.29),
                         f"{arrow} {abs(mom):.0f}% MoM",
                         font_size=Pt(9), font_color=color, bold=True)

        if label == "Quarter to Date" and value != "N/A":
            _add_textbox(slide, x + Inches(0.18), card_y + Inches(1.06), card_w - Inches(0.2), Inches(0.25),
                         "Cumulative Q1 revenue",
                         font_size=Pt(7), font_color=GRAY)

    # ── Gauge Charts ──
    gauge_card = _add_card(slide, Inches(0.30), Inches(4.15), Inches(7.89), Inches(2.20))

    gauges = [
        ("% of Net Revenue Goal", data.pct_net_revenue_goal),
        ("% of AOV Goal", data.pct_aov_goal),
        ("Utilization", data.utilization_rate),
        ("Rebooking Rate", data.rebooking_rate),
        ("Retention (180D)", data.retention_180d),
    ]

    gauge_w = Inches(1.12)
    gauge_spacing = Inches(1.43)
    gauge_x_start = Inches(0.55)
    gauge_y = Inches(4.25)
    label_y = Inches(5.55)

    for i, (label, val) in enumerate(gauges):
        x = gauge_x_start + i * gauge_spacing

        # Generate and insert gauge image
        png_bytes = make_gauge_png(val, size=300, line_width=28)
        img_stream = BytesIO(png_bytes)
        slide.shapes.add_picture(img_stream, x, gauge_y, gauge_w, gauge_w)

        # Percentage text centered on gauge
        pct_text = f"{val * 100:.0f}%" if val == int(val * 100) / 100 else f"{val * 100:.1f}%"
        if val * 100 == int(val * 100):
            pct_text = f"{int(val * 100)}%"
        _add_textbox(slide, x + Inches(0.15), gauge_y + Inches(0.35), Inches(0.82), Inches(0.41),
                     pct_text, font_size=Pt(16), font_color=PURPLE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Label below gauge
        _add_textbox(slide, x - Inches(0.15), label_y, Inches(1.45), Inches(0.30),
                     label, font_size=Pt(7.5), font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Memberships ──
    mem_y = Inches(6.55)
    _add_card(slide, Inches(0.30), mem_y, Inches(4.05), Inches(1.55))
    _add_textbox(slide, Inches(0.50), mem_y + Inches(0.13), Inches(2.5), Inches(0.43),
                 "Memberships", font_name=FONT_HEADING, font_size=Pt(12.5), font_color=PURPLE, bold=True)

    mem_items = [
        (data.memberships_active, "Active"),
        (data.memberships_new, "New"),
        (data.memberships_cancelled, "Cancelled"),
        (_fmt_dollar(data.mrr), "MRR"),
    ]
    mem_x_start = Inches(0.35)
    mem_col_w = Inches(0.95)
    for i, (val, label) in enumerate(mem_items):
        x = mem_x_start + i * mem_col_w
        val_size = Pt(17) if len(str(val)) > 4 else Pt(21)
        _add_textbox(slide, x, mem_y + Inches(0.64), Inches(1.08), Inches(0.42),
                     str(val), font_size=val_size, font_color=PURPLE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x, mem_y + Inches(1.06), Inches(1.08), Inches(0.20),
                     label, font_size=Pt(9), font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Client Mix ──
    _add_card(slide, Inches(4.54), mem_y, Inches(3.65), Inches(1.55))
    _add_textbox(slide, Inches(4.70), mem_y + Inches(0.12), Inches(2.0), Inches(0.42),
                 "Client Mix", font_name=FONT_HEADING, font_size=Pt(14.5), font_color=PURPLE, bold=True)

    # Stacked bar
    bar_y = mem_y + Inches(0.63)
    total = data.total_clients or 1
    existing_pct = data.existing_clients / total
    bar_total_w = Inches(3.40)

    existing_w = int(bar_total_w * existing_pct)
    new_w = int(bar_total_w) - existing_w

    if existing_w > 0:
        bar_existing = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               Inches(4.70), bar_y, Emu(existing_w), Inches(0.28))
        bar_existing.fill.solid()
        bar_existing.fill.fore_color.rgb = PURPLE
        bar_existing.line.fill.background()

    if new_w > 0:
        bar_new = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(4.70) + Emu(existing_w), bar_y, Emu(new_w), Inches(0.28))
        bar_new.fill.solid()
        bar_new.fill.fore_color.rgb = GOLD
        bar_new.line.fill.background()

    # Labels
    _add_textbox(slide, Inches(4.70), bar_y + Inches(0.40), Inches(1.5), Inches(0.23),
                 f"Existing {data.existing_client_pct:.0f}%",
                 font_size=Pt(8.5), font_color=LIGHT_GRAY)
    _add_textbox(slide, Inches(7.0), bar_y + Inches(0.40), Inches(1.2), Inches(0.23),
                 f"New {data.new_client_pct:.0f}%",
                 font_size=Pt(8.5), font_color=LIGHT_GRAY)

    # ── Ratings & Reviews ──
    review_y = Inches(8.29)
    _add_card(slide, Inches(0.30), review_y, Inches(7.89), Inches(1.82))
    _add_textbox(slide, Inches(0.47), review_y - Inches(0.04), Inches(3.28), Inches(0.38),
                 "Ratings & Reviews",
                 font_name=FONT_HEADING, font_size=Pt(12.5), font_color=GOLD, bold=True)

    # Table headers
    col_starts = [Inches(0.40), Inches(1.77), Inches(3.46), Inches(4.99), Inches(6.52)]
    col_widths = [Inches(1.38), Inches(1.68), Inches(1.53), Inches(1.53), Inches(1.53)]
    headers = ["Platform", "New Reviews\nthis Month", "Avg. New\nRating", "Total\nReviews", "Overall\nRating"]
    header_y = review_y + Inches(0.38)

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        hdr_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, header_y, cw, Inches(0.42))
        hdr_shape.fill.solid()
        hdr_shape.fill.fore_color.rgb = TABLE_HEADER_BG
        hdr_shape.line.fill.background()

        _add_textbox(slide, cx, header_y, cw, Inches(0.42),
                     hdr, font_size=Pt(7.5), font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Review rows
    if data.reviews:
        for ri, rev in enumerate(data.reviews):
            row_y = header_y + Inches(0.42) + ri * Inches(0.37)
            vals = [
                rev.platform,
                str(rev.new_reviews) if rev.new_reviews is not None else "-",
                str(rev.avg_new_rating) if rev.avg_new_rating is not None else "-",
                str(rev.total_reviews) if rev.total_reviews is not None else "-",
                str(rev.overall_rating) if rev.overall_rating is not None else "-",
            ]
            for j, (v, cx, cw) in enumerate(zip(vals, col_starts, col_widths)):
                _add_textbox(slide, cx, row_y, cw, Inches(0.37),
                             v, font_size=Pt(8), font_color=DARK_TEXT,
                             alignment=PP_ALIGN.CENTER)
    else:
        _add_textbox(slide, Inches(1.5), header_y + Inches(0.55), Inches(5.0), Inches(0.37),
                     "Data not available", font_size=Pt(9), font_color=GRAY,
                     alignment=PP_ALIGN.CENTER)

    _add_copyright(slide)


def build_slide_3_revenue_staff(prs, data: MBRData):
    """Slide 3: Revenue Breakdown + Staff Performance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    # ── Revenue Breakdown ──
    _add_textbox(slide, Inches(0.57), Inches(0.55), Inches(3.5), Inches(0.40),
                 "Revenue Breakdown", font_name=FONT_HEADING, font_size=Pt(16),
                 font_color=PINK_ACCENT, bold=True)

    _add_card(slide, Inches(0.49), Inches(1.00), Inches(7.60), Inches(3.15))

    _add_textbox(slide, Inches(0.65), Inches(1.10), Inches(3.0), Inches(0.37),
                 "Gross Revenue Sources", font_name=FONT_HEADING, font_size=Pt(11),
                 font_color=PURPLE, bold=True)

    # Revenue bars
    rev_items = [
        ("Service Revenue", data.service_revenue),
        ("Prepayment Revenue", data.prepayment_revenue),
        ("Membership Sales", data.membership_sales),
        ("Custom Items", data.custom_items),
        ("Retail Revenue", data.retail_revenue),
    ]

    max_rev = max(r[1] for r in rev_items) if rev_items else 1
    bar_x_label = Inches(0.58)
    bar_x_start = Inches(2.31)
    bar_max_w = Inches(4.49)
    bar_h = Inches(0.21)

    for i, (label, val) in enumerate(rev_items):
        y = Inches(1.55) + i * Inches(0.36)

        # Label
        _add_textbox(slide, bar_x_label, y - Inches(0.03), Inches(1.7), Inches(0.26),
                     label, font_size=Pt(8), font_color=PURPLE)

        # Background track
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        bar_x_start, y, bar_max_w, bar_h)
        track.fill.solid()
        track.fill.fore_color.rgb = TRACK_COLOR
        track.line.fill.background()

        # Filled bar
        if val > 0 and max_rev > 0:
            fill_w = max(int(bar_max_w * (val / max_rev)), Emu(Inches(0.05)))
            bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               bar_x_start, y, fill_w, bar_h)
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = PURPLE
            bar_fill.line.fill.background()

        # Value
        _add_textbox(slide, Inches(6.80), y - Inches(0.03), Inches(1.0), Inches(0.26),
                     _fmt_dollar(val), font_size=Pt(8), font_color=PURPLE,
                     alignment=PP_ALIGN.RIGHT)

    # Total line
    total_y = Inches(1.55) + 5 * Inches(0.36) + Inches(0.08)
    _add_textbox(slide, bar_x_label, total_y, Inches(1.7), Inches(0.22),
                 "Total Gross", font_size=Pt(8), font_color=PURPLE, bold=True)
    _add_textbox(slide, Inches(6.80), total_y, Inches(1.0), Inches(0.22),
                 _fmt_dollar(data.total_gross), font_size=Pt(8), font_color=PURPLE,
                 bold=True, alignment=PP_ALIGN.RIGHT)

    # ── Adjustments (3 cards) ──
    adj_y = Inches(3.80)
    adj_items = [
        ("Discounts", data.discounts, True),
        ("Redemptions", data.redemptions, True),
        ("Client Fees", data.client_fees, False),
    ]
    adj_card_w = Inches(2.50)

    for i, (label, val, is_negative) in enumerate(adj_items):
        x = Inches(0.51) + i * Inches(2.69)
        _add_card(slide, x, adj_y, adj_card_w, Inches(0.75))
        _add_textbox(slide, x + Inches(0.13), adj_y + Inches(0.02), Inches(2.2), Inches(0.30),
                     label, font_name=FONT_HEADING, font_size=Pt(10), font_color=PURPLE, bold=True)

        prefix = "-" if is_negative and val > 0 else "+"
        color = RED if is_negative else GREEN
        display_val = f"{prefix}{_fmt_dollar(val)}"
        if val == 0:
            display_val = "$0"
            color = GRAY
        _add_textbox(slide, x + Inches(0.13), adj_y + Inches(0.30), Inches(2.2), Inches(0.37),
                     display_val, font_size=Pt(14), font_color=color, bold=True)

    # ── Staff Performance ──
    staff_section_y = Inches(4.85)
    _add_textbox(slide, Inches(0.51), staff_section_y, Inches(3.5), Inches(0.40),
                 "Staff Performance", font_name=FONT_HEADING, font_size=Pt(16),
                 font_color=PINK_ACCENT, bold=True)

    staff = data.staff
    if not staff:
        _add_textbox(slide, Inches(1.0), staff_section_y + Inches(0.6), Inches(6.0), Inches(0.5),
                     "No staff data available", font_size=Pt(11), font_color=GRAY,
                     alignment=PP_ALIGN.CENTER)
        _add_copyright(slide)
        return

    # Dynamic layout: 2 columns, adaptive card height based on staff count
    num_staff = len(staff)
    cols = 2
    card_w = Inches(3.65)
    x_starts = [Inches(0.57), Inches(4.47)]
    start_y = staff_section_y + Inches(0.55)
    num_rows = (num_staff + cols - 1) // cols

    # Calculate card height to fit within page (leave room for contribution bar)
    available_h = Inches(10.0) - start_y - Inches(0.70)  # page - start - contrib bar
    row_spacing = Inches(0.08)
    card_h = min(Inches(2.18), (available_h - (num_rows - 1) * row_spacing) / num_rows)
    # Determine if we need compact mode
    compact = card_h < Inches(1.80)

    total_staff_rev = sum(s.net_revenue for s in staff)

    for si, s in enumerate(staff):
        col = si % cols
        row = si // cols
        x = x_starts[col]
        y = start_y + row * (card_h + row_spacing)

        # Card background
        _add_card(slide, x, y, card_w, card_h)

        # Top accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PURPLE
        accent.line.fill.background()

        # Initials circle
        circ_size = Inches(0.36) if compact else Inches(0.46)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.10), y + Inches(0.12),
                                         circ_size, circ_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = PURPLE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = s.initials
        run.font.size = Pt(9) if compact else Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True
        tf.paragraphs[0].space_before = Pt(2)

        # Name and revenue
        name_x = x + Inches(0.52) if compact else x + Inches(0.66)
        _add_textbox(slide, name_x, y + Inches(0.10), card_w - Inches(0.8), Inches(0.22),
                     s.name, font_size=Pt(8), font_color=PURPLE, bold=True)
        _add_textbox(slide, name_x, y + Inches(0.30), card_w - Inches(0.8), Inches(0.16),
                     f"{_fmt_dollar(s.net_revenue)} net revenue",
                     font_size=Pt(6.5), font_color=GRAY)

        # Metric boxes (AOV, Utilization, Rebooking)
        metric_y = y + Inches(0.52)
        inner_w = (card_w - Inches(0.24)) / 3
        metric_h = Inches(0.46) if compact else Inches(0.63)
        staff_metrics = [
            ("AOV", _fmt_dollar(s.aov)),
            ("Utilization", _fmt_pct(s.utilization) if s.utilization else "N/A"),
            ("Rebooking", _fmt_pct(s.rebooking_rate) if s.rebooking_rate else "N/A"),
        ]
        for mi, (mlabel, mval) in enumerate(staff_metrics):
            mx = x + Inches(0.12) + mi * inner_w
            _add_card(slide, mx, metric_y, inner_w - Inches(0.04), metric_h,
                      fill_color=RGBColor(0xF8, 0xF5, 0xF0))
            _add_textbox(slide, mx, metric_y + Inches(0.04), inner_w - Inches(0.04), Inches(0.15),
                         mlabel, font_size=Pt(6), font_color=GRAY, alignment=PP_ALIGN.CENTER)
            _add_textbox(slide, mx, metric_y + Inches(0.20), inner_w - Inches(0.04), Inches(0.25),
                         mval, font_size=Pt(9), font_color=PURPLE, bold=True,
                         alignment=PP_ALIGN.CENTER)

        # Service + Retail Revenue — inline layout in compact mode
        bottom_y = metric_y + metric_h + Inches(0.05)
        half_w = (card_w - Inches(0.24)) / 2
        if compact:
            # Single row: "Svc Rev: $X | Retail: $Y"
            for bi, (blabel, bval) in enumerate([
                ("Svc Rev", _fmt_dollar(s.service_revenue)),
                ("Retail", _fmt_dollar(s.retail_revenue)),
            ]):
                bx = x + Inches(0.12) + bi * half_w
                _add_card(slide, bx, bottom_y, half_w - Inches(0.04), Inches(0.32),
                          fill_color=RGBColor(0xF8, 0xF5, 0xF0))
                _add_textbox(slide, bx + Inches(0.06), bottom_y + Inches(0.02),
                             half_w - Inches(0.12), Inches(0.13),
                             blabel, font_size=Pt(5), font_color=GRAY)
                _add_textbox(slide, bx + Inches(0.06), bottom_y + Inches(0.14),
                             half_w - Inches(0.12), Inches(0.16),
                             bval, font_size=Pt(8), font_color=PURPLE, bold=True)
        else:
            bottom_h = Inches(0.55)
            for bi, (blabel, bval) in enumerate([
                ("Service Revenue", _fmt_dollar(s.service_revenue)),
                ("Retail Revenue", _fmt_dollar(s.retail_revenue)),
            ]):
                bx = x + Inches(0.12) + bi * half_w
                _add_card(slide, bx, bottom_y, half_w - Inches(0.04), bottom_h,
                          fill_color=RGBColor(0xF8, 0xF5, 0xF0))
                _add_textbox(slide, bx + Inches(0.06), bottom_y + Inches(0.03),
                             half_w - Inches(0.12), Inches(0.13),
                             blabel, font_size=Pt(5.5), font_color=GRAY)
                _add_textbox(slide, bx + Inches(0.06), bottom_y + Inches(0.17),
                             half_w - Inches(0.12), Inches(0.22),
                             bval, font_size=Pt(9), font_color=PURPLE, bold=True)

    # ── Revenue Contribution bar ──
    contrib_y = start_y + ((num_staff + cols - 1) // cols) * (card_h + Inches(0.15)) + Inches(0.10)
    # Only show if it fits on the page
    if contrib_y + Inches(0.60) < Inches(10.5):
        _add_textbox(slide, Inches(0.57), contrib_y, Inches(3.5), Inches(0.19),
                     "Revenue Contribution", font_size=Pt(8), font_color=PURPLE, bold=True)

        bar_y = contrib_y + Inches(0.28)
        bar_w = Inches(7.50)
        colors = [PURPLE, GOLD, PINK_ACCENT, RGBColor(0x7B, 0x4D, 0x8E),
                  RGBColor(0xD4, 0xA7, 0x3A), RGBColor(0x9C, 0x6D, 0xA0)]

        x_offset = Inches(0.57)
        for si, s in enumerate(staff):
            pct = s.net_revenue / total_staff_rev if total_staff_rev > 0 else 0
            segment_w = int(bar_w * pct)
            if segment_w > 0:
                seg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              x_offset, bar_y, segment_w, Inches(0.35))
                seg.fill.solid()
                seg.fill.fore_color.rgb = colors[si % len(colors)]
                seg.line.fill.background()

                # Label inside bar if wide enough
                short_name = s.name.split()[0]
                seg_label = f"{short_name}  \u00b7  {pct * 100:.0f}%"
                if segment_w > Inches(0.8):
                    _add_textbox(slide, x_offset + Emu(Inches(0.04)), bar_y + Emu(Inches(0.03)),
                                 segment_w - Emu(Inches(0.08)), Inches(0.29),
                                 seg_label, font_size=Pt(7), font_color=WHITE, bold=True)

                x_offset += segment_w

    _add_copyright(slide)


def build_slide_4_services_marketing(prs, data: MBRData):
    """Slide 4: Service Mix + Marketing Performance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    # ── Service Mix Title ──
    _add_textbox(slide, Inches(0.57), Inches(0.55), Inches(3.5), Inches(0.40),
                 "Service Mix", font_name=FONT_HEADING, font_size=Pt(16),
                 font_color=PINK_ACCENT, bold=True)

    # Top 3 services highlight cards
    svc_card_h = Inches(1.80) + len(data.services[:10]) * Inches(0.22) + Inches(0.30)
    _add_card(slide, Inches(0.43), Inches(1.00), Inches(7.60), svc_card_h)

    services = data.services[:10]  # Show up to 10
    top3 = services[:3]
    top3_card_w = Inches(2.19)

    for i, svc in enumerate(top3):
        x = Inches(0.91) + i * Inches(2.43)
        y = Inches(1.15)
        _add_card(slide, x, y, top3_card_w, Inches(1.12), fill_color=RGBColor(0xF8, 0xF5, 0xF0))

        # Accent top bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, top3_card_w, Inches(0.05))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PURPLE
        accent.line.fill.background()

        ordinal = ["1st", "2nd", "3rd"][i]
        _add_textbox(slide, x, y + Inches(0.12), top3_card_w, Inches(0.21),
                     ordinal, font_size=Pt(9), font_color=GRAY, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.12), y + Inches(0.30), top3_card_w - Inches(0.24), Inches(0.23),
                     svc.name, font_size=Pt(9), font_color=PURPLE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.58), top3_card_w, Inches(0.33),
                     f"{svc.pct_of_total:.1f}%", font_size=Pt(18), font_color=PURPLE,
                     bold=True, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.90), top3_card_w, Inches(0.17),
                     "of total revenue", font_size=Pt(7), font_color=GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Full service bars
    _add_textbox(slide, Inches(0.67), Inches(2.60), Inches(4.5), Inches(0.22),
                 "% of Total Revenue by Service", font_size=Pt(8), font_color=PURPLE, bold=True)

    max_pct = services[0].pct_of_total if services else 1
    bar_x_start = Inches(2.10)
    bar_max_w = Inches(5.49)
    bar_h = Inches(0.13)
    row_h = Inches(0.22)

    for i, svc in enumerate(services):
        y = Inches(2.95) + i * row_h

        # Truncate long service names
        name = svc.name if len(svc.name) <= 25 else svc.name[:23] + ".."
        _add_textbox(slide, Inches(0.51), y - Inches(0.02), Inches(1.55), Inches(0.18),
                     name, font_size=Pt(5.5), font_color=PURPLE)

        # Background track
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        bar_x_start, y, bar_max_w, bar_h)
        track.fill.solid()
        track.fill.fore_color.rgb = TRACK_COLOR
        track.line.fill.background()

        # Fill bar
        if svc.pct_of_total > 0:
            fill_w = max(int(bar_max_w * (svc.pct_of_total / max_pct)), Emu(Inches(0.03)))
            bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               bar_x_start, y, fill_w, bar_h)
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = PURPLE
            bar_fill.line.fill.background()

        # Percentage
        _add_textbox(slide, Inches(7.66), y - Inches(0.02), Inches(0.50), Inches(0.17),
                     f"{svc.pct_of_total:.1f}%", font_size=Pt(6), font_color=PURPLE,
                     alignment=PP_ALIGN.RIGHT)

    # ── Marketing Performance ──
    num_services = len(services)
    mkt_y = Inches(2.95) + num_services * row_h + Inches(0.50)
    _add_textbox(slide, Inches(0.59), mkt_y, Inches(3.0), Inches(0.41),
                 "Marketing Performance", font_name=FONT_HEADING, font_size=Pt(14),
                 font_color=PINK_ACCENT, bold=True)

    _add_card(slide, Inches(0.31), mkt_y + Inches(0.55), Inches(7.89), Inches(4.50))

    if data.marketing:
        mkt = data.marketing
        # ROI headline
        roi_val = mkt.revenue / mkt.ad_spend if mkt.ad_spend > 0 else 0
        _add_card(slide, Inches(0.74), mkt_y + Inches(0.70), Inches(7.08), Inches(0.59),
                  fill_color=RGBColor(0xF8, 0xF5, 0xF0))
        roi_text = f"For every $1 you spend on this campaign, you generate"
        _add_textbox(slide, Inches(1.03), mkt_y + Inches(0.72), Inches(4.5), Inches(0.30),
                     roi_text, font_size=Pt(8), font_color=PURPLE)
        _add_textbox(slide, Inches(1.03), mkt_y + Inches(1.00), Inches(1.5), Inches(0.30),
                     f"${roi_val:.2f}", font_size=Pt(14), font_color=PURPLE, bold=True)
        _add_textbox(slide, Inches(2.10), mkt_y + Inches(1.00), Inches(4.0), Inches(0.30),
                     "from new patients on their first visit",
                     font_size=Pt(8), font_color=PURPLE)

        # Funnel cards
        funnel = [
            ("Ad Spend", _fmt_dollar(mkt.ad_spend), "Monthly Budget"),
            ("Leads", str(mkt.leads), "New Patient Leads"),
            ("Booked", str(mkt.booked), "# of Booked Appointments"),
            ("Completed", str(mkt.completed), "# of Completed Appointments"),
            ("Revenue", _fmt_dollar(mkt.revenue), "First-visit Revenue"),
        ]
        for fi, (flabel, fval, fsub) in enumerate(funnel):
            fx = Inches(0.60) + fi * Inches(1.52)
            fy = mkt_y + Inches(1.55)
            _add_card(slide, fx, fy, Inches(1.37), Inches(0.92),
                      fill_color=RGBColor(0xF8, 0xF5, 0xF0))
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, fx, fy, Inches(1.37), Inches(0.03))
            accent.fill.solid()
            accent.fill.fore_color.rgb = PURPLE
            accent.line.fill.background()
            _add_textbox(slide, fx + Inches(0.08), fy + Inches(0.10), Inches(1.23), Inches(0.17),
                         flabel, font_size=Pt(7), font_color=GRAY, alignment=PP_ALIGN.CENTER)
            _add_textbox(slide, fx, fy + Inches(0.30), Inches(1.37), Inches(0.24),
                         fval, font_size=Pt(12), font_color=PURPLE, bold=True,
                         alignment=PP_ALIGN.CENTER)
            _add_textbox(slide, fx, fy + Inches(0.68), Inches(1.37), Inches(0.12),
                         fsub, font_size=Pt(5), font_color=GRAY, alignment=PP_ALIGN.CENTER)

        # Key Metrics
        km_y = mkt_y + Inches(2.70)
        _add_textbox(slide, Inches(0.67), km_y, Inches(2.0), Inches(0.17),
                     "Key Metrics", font_size=Pt(8), font_color=PURPLE, bold=True)

        key_metrics = [
            ("First-visit ROI", f"{mkt.first_visit_roi:.2f}x" if mkt.first_visit_roi else "N/A",
             "Goal: 3x", mkt.first_visit_roi and mkt.first_visit_roi >= 3),
            ("Lead \u2192 Booking Rate",
             f"{mkt.lead_to_booking_rate:.2f}%" if mkt.lead_to_booking_rate else "N/A",
             "Goal: 15%", mkt.lead_to_booking_rate and mkt.lead_to_booking_rate >= 15),
            ("First-Visit AOV", _fmt_dollar(mkt.first_visit_aov) if mkt.first_visit_aov else "N/A",
             "Goal: $575", mkt.first_visit_aov and mkt.first_visit_aov >= 575),
        ]
        for ki, (klabel, kval, kgoal, on_track) in enumerate(key_metrics):
            ky = km_y + Inches(0.28) + ki * Inches(0.46)
            _add_textbox(slide, Inches(0.69), ky, Inches(1.5), Inches(0.21),
                         klabel, font_size=Pt(8), font_color=PURPLE)
            _add_textbox(slide, Inches(0.70), ky + Inches(0.19), Inches(1.5), Inches(0.14),
                         kgoal, font_size=Pt(6), font_color=GRAY)
            _add_textbox(slide, Inches(2.36), ky - Inches(0.04), Inches(1.0), Inches(0.30),
                         kval, font_size=Pt(11), font_color=PURPLE, bold=True)

            # Status badge
            badge_text = "Meets Target" if on_track else "Below Target"
            badge_color = GREEN if on_track else GOLD
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(3.37), ky, Inches(0.70), Inches(0.20))
            badge.fill.solid()
            badge.fill.fore_color.rgb = badge_color
            badge.line.fill.background()
            _add_textbox(slide, Inches(3.37), ky, Inches(0.70), Inches(0.20),
                         badge_text, font_size=Pt(5.5), font_color=WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)

        # Next Steps
        _add_textbox(slide, Inches(4.49), km_y, Inches(2.0), Inches(0.17),
                     "Next Steps", font_size=Pt(8), font_color=PURPLE, bold=True)
        for ni, step in enumerate(mkt.next_steps[:3]):
            ny = km_y + Inches(0.30) + ni * Inches(0.46)
            _add_card(slide, Inches(4.51), ny, Inches(3.48), Inches(0.40),
                      fill_color=RGBColor(0xF8, 0xF5, 0xF0))
            _add_textbox(slide, Inches(4.58), ny + Inches(0.04), Inches(3.30), Inches(0.32),
                         step, font_size=Pt(6), font_color=PURPLE)
    else:
        _add_textbox(slide, Inches(2.0), mkt_y + Inches(2.0), Inches(4.5), Inches(0.5),
                     "Marketing data not available for this period.",
                     font_size=Pt(11), font_color=GRAY, alignment=PP_ALIGN.CENTER)

    _add_copyright(slide)


def build_slide_5_takeaways(prs, data: MBRData):
    """Slide 5: Takeaways & Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    _add_textbox(slide, Inches(0.71), Inches(0.55), Inches(6.5), Inches(0.44),
                 "Takeaways & Recommendations",
                 font_name=FONT_HEADING, font_size=Pt(18), font_color=PINK_ACCENT, bold=True)

    # ── Assessment Section ──
    _add_card(slide, Inches(0.43), Inches(1.15), Inches(7.60), Inches(4.20))
    _add_textbox(slide, Inches(0.89), Inches(1.25), Inches(4.5), Inches(0.22),
                 f"{data.month_name} Assessment",
                 font_size=Pt(9), font_color=PURPLE, bold=True)

    tag_colors = {
        "STRENGTH": GREEN,
        "OPPORTUNITY": GOLD,
        "WARNING": RED,
    }

    assessments = data.assessments or []
    for i, assess in enumerate(assessments[:5]):
        y = Inches(1.60) + i * Inches(0.60)

        # Card background
        _add_card(slide, Inches(0.84), y, Inches(6.82), Inches(0.53),
                  fill_color=RGBColor(0xF8, 0xF5, 0xF0))

        # Left accent bar
        tag = assess.get("tag", "OPPORTUNITY")
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.84), y, Inches(0.05), Inches(0.53))
        accent.fill.solid()
        accent.fill.fore_color.rgb = tag_colors.get(tag, GOLD)
        accent.line.fill.background()

        # Tag badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(1.05), y + Inches(0.07), Inches(0.70), Inches(0.15))
        badge.fill.solid()
        badge.fill.fore_color.rgb = tag_colors.get(tag, GOLD)
        badge.line.fill.background()
        _add_textbox(slide, Inches(1.05), y + Inches(0.06), Inches(0.70), Inches(0.16),
                     tag, font_size=Pt(5.5), font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        _add_textbox(slide, Inches(1.84), y + Inches(0.05), Inches(5.5), Inches(0.18),
                     assess.get("title", ""), font_size=Pt(9), font_color=PURPLE, bold=True)

        # Text
        _add_textbox(slide, Inches(1.05), y + Inches(0.24), Inches(6.36), Inches(0.26),
                     assess.get("text", ""), font_size=Pt(7), font_color=DARK_TEXT)

    # ── PSM Feedback Section ──
    psm_y = Inches(5.60)
    _add_card(slide, Inches(0.39), psm_y, Inches(7.60), Inches(4.60))

    _add_textbox(slide, Inches(0.89), psm_y + Inches(0.20), Inches(6.5), Inches(0.38),
                 "Feedback from your Practice Success Manager",
                 font_size=Pt(10), font_color=PURPLE, bold=True)

    if data.psm_feedback:
        _add_textbox(slide, Inches(0.89), psm_y + Inches(0.65), Inches(6.36), Inches(3.70),
                     data.psm_feedback, font_size=Pt(8), font_color=DARK_TEXT)

    _add_copyright(slide)


def build_slide_6_brand_bank(prs, data: MBRData, brand_bank_path=None):
    """Slide 6: Brand Bank (static/optional)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    _add_card(slide, Inches(0.51), Inches(0.60), Inches(7.50), Inches(9.40))

    next_month_num = data.month % 12 + 1
    next_month_name = __import__('calendar').month_name[next_month_num]

    _add_textbox(slide, Inches(0.51), Inches(0.70), Inches(3.0), Inches(0.61),
                 f"{next_month_name}\nBrand Bank",
                 font_name=FONT_HEADING, font_size=Pt(16), font_color=PURPLE, bold=True)

    if brand_bank_path and os.path.exists(brand_bank_path):
        slide.shapes.add_picture(brand_bank_path, Inches(0.75), Inches(1.60),
                                  width=Inches(7.0))
    else:
        _add_textbox(slide, Inches(1.5), Inches(4.5), Inches(5.5), Inches(1.0),
                     "Brand Bank content will be provided by the Marketing team.",
                     font_size=Pt(14), font_color=GRAY, alignment=PP_ALIGN.CENTER)

    _add_copyright(slide)


def build_slide_7_launches(prs, data: MBRData):
    """Slide 7: Moxie Suite Launches (static/optional)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    _add_textbox(slide, Inches(0.21), Inches(0.70), Inches(8.0), Inches(0.43),
                 "Moxie Suite Launches | New and Exciting",
                 font_name=FONT_HEADING, font_size=Pt(18), font_color=PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Placeholder cards for launches
    launches = [
        ("No Provider Preference\nOnline Booking",
         "Let clients book the first available provider for a service, so they can snag the soonest appointment."),
        ("Calendar Flow",
         "Build schedules that match how treatments actually run day to day with custom durations and buffers."),
        ("Product Presets",
         "Save your most common product formulas right inside each service for smoother clinical days."),
        ("Automatic Retail\nTax Updates",
         "Retail tax rates now update automatically based on your product type and location."),
    ]

    positions = [
        (Inches(0.70), Inches(1.50)), (Inches(4.49), Inches(1.50)),
        (Inches(0.61), Inches(5.50)), (Inches(4.47), Inches(5.50)),
    ]

    for i, ((title, desc), (x, y)) in enumerate(zip(launches, positions)):
        _add_card(slide, x, y, Inches(3.32), Inches(3.50))
        _add_textbox(slide, x + Inches(0.30), y + Inches(0.30), Inches(2.70), Inches(0.50),
                     title, font_name=FONT_HEADING, font_size=Pt(12), font_color=PURPLE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.20), y + Inches(0.90), Inches(2.90), Inches(2.30),
                     desc, font_size=Pt(9), font_color=DARK_TEXT)

    _add_copyright(slide)


def build_slide_8_partnerships(prs, data: MBRData):
    """Slide 8: Device Partnerships (static)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header_bar(slide, data.practice_name, f"{data.month_name} {data.year}")

    _add_textbox(slide, Inches(0.33), Inches(2.80), Inches(7.84), Inches(0.50),
                 "Considering adding a new device to your practice?",
                 font_name=FONT_HEADING, font_size=Pt(20), font_color=PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.0), Inches(3.80), Inches(8.50), Inches(0.71))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PURPLE
    banner.line.fill.background()
    _add_textbox(slide, Inches(0.5), Inches(3.82), Inches(7.50), Inches(0.65),
                 "Before you make a decision, connect with your PSM to explore\nMoxie's exclusive discounts with top companies!",
                 font_size=Pt(11), font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # Partner logo cards
    partners = ["Cynosure", "BTL", "Sciton", "And more!"]
    for i, name in enumerate(partners):
        x = Inches(0.54) + i * Inches(1.96)
        _add_card(slide, x, Inches(5.00), Inches(1.78), Inches(2.46))

        # Try to use extracted partner logos
        logo_path = ASSETS_DIR / f"partner_{i}.png"
        if logo_path.exists() and i < 3:
            slide.shapes.add_picture(str(logo_path), x + Inches(0.15), Inches(5.50),
                                      width=Inches(1.48))
        else:
            _add_textbox(slide, x, Inches(5.80), Inches(1.78), Inches(0.50),
                         name, font_name=FONT_HEADING, font_size=Pt(14),
                         font_color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_copyright(slide)


def build_mbr(data: MBRData, output_path: str, brand_bank_path: str = None):
    """Build the complete MBR presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_cover(prs, data)
    build_slide_2_snapshot(prs, data)
    build_slide_3_revenue_staff(prs, data)
    build_slide_4_services_marketing(prs, data)
    build_slide_5_takeaways(prs, data)
    build_slide_6_brand_bank(prs, data, brand_bank_path)
    build_slide_7_launches(prs, data)
    build_slide_8_partnerships(prs, data)

    prs.save(output_path)
    return output_path
